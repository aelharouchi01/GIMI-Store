# -*- coding: utf-8 -*-
from pptx import Presentation
import os

prs = Presentation(r'C:\Users\aelha\OneDrive\Desktop\GIMI_Product_Catalogue_V64.pptx')

# Dump all 64 slides to file, then search
out = []
for i, slide in enumerate(prs.slides, 1):
    text = ''
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    text += run.text + ' '
                text += '\n'
    out.append(f'=== SLIDE {i} ===\n{text}\n')

with open(r'C:\Users\aelha\Downloads\deck_full.txt', 'w', encoding='utf-8') as f:
    f.write('\n'.join(out))

print(f'Done. {len(prs.slides)} slides written.')
