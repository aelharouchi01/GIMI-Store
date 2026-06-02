"""Build GIMI Product Catalogue V51 — 62 slides aligned to website 8-section structure."""
import json, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ===== Brand =====
TEAL       = RGBColor(0x2A, 0xAC, 0xB5)
TEAL_DARK  = RGBColor(0x1E, 0x8A, 0x93)
TEAL_VDARK = RGBColor(0x10, 0x4E, 0x60)
YELLOW     = RGBColor(0xC9, 0xD9, 0x40)
NAVY       = RGBColor(0x1B, 0x1F, 0x3B)
NAVY_DARK  = RGBColor(0x13, 0x16, 0x2C)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY       = RGBColor(0x6B, 0x72, 0x80)
GRAY_LIGHT = RGBColor(0xF0, 0xF4, 0xF6)
TEXT       = RGBColor(0x1F, 0x29, 0x37)
PINK       = RGBColor(0xE1, 0x5A, 0x78)
ORANGE     = RGBColor(0xF0, 0x87, 0x1E)

FONT_TITLE = 'Calibri'
FONT_BODY  = 'Calibri'

LOGO_DARK  = r'C:\Users\aelha\Downloads\gimi-store-assets\gimi-logo.png'
LOGO_WHITE = r'C:\Users\aelha\Downloads\gimi-store-assets\gimi-logo-white.png'

# ===== Load products =====
products = json.load(open(r'C:\Users\aelha\Downloads\products.json'))

# Apply modality defaults (mirroring the mk() helper in the website JS)
def default_modality(p):
    if p.get('modality'): return p['modality']
    if p.get('cluster') in ('cip-mc', 'ccio-mc'): return 'Live Virtual Cohort'
    if p['track'] == 'books': return 'Reference (digital)'
    if p['track'] == 'memb' and p['id'] == 'thinktank': return 'Monthly virtual sessions'
    if p['track'] == 'memb': return 'Annual access'
    if p.get('cluster') in ('indiv-ass', 'org-ass'): return 'Online assessment'
    if p.get('cluster') == 'software': return 'Software platform'
    if p.get('cluster') == 'engage': return 'Custom engagement'
    return 'Online Self-Paced'

for p in products:
    p['modality'] = default_modality(p)

# Track labels (8 sections)
TRACKS = [
    ('icore', 'Innovation Core',     'The flagship innovation management pathway'),
    ('ielec', 'Innovation Elective', 'Specialist tracks: Design Thinking, ISO, AI, Tech, Longevity, Audit'),
    ('ff',    'Future Foresight',    'Foresight professional and trainer certifications'),
    ('lead',  'Leadership Core',     'Leadership-oriented certifications'),
    ('cons',  'Consulting Core',     'The Management Consulting Institute pathway'),
    ('books', 'Books & Guides',      'Published books, guides, and reference materials'),
    ('tools', 'Programs & Tools',    'Assessments, software, programs, and engagements'),
    ('memb',  'Memberships',         'Annual access tiers for individuals, companies, and academia'),
]

CLUSTERS = {
    'cip': 'CIP series', 'cip-mc': 'CIP-Masterclass', 'ccio': 'CCIO series',
    'ccio-mc': 'CCIO-Masterclass', 'impact': 'GIMI Impact', 'icore-train': 'Train the Trainer',
    'ff-series': 'Foresight series', 'ff-train': 'Train the Trainer',
    'imbok': 'IMBOK', 'ffbok': 'Future Foresight BOK', 'iso-bk': 'ISO Standards',
    'stories': 'Innovation Stories', 'tool-bk': 'Practitioner Tools', 'journal': 'Journal',
    'indiv-ass': 'Individual Assessments', 'org-ass': 'Organisational Assessments',
    'software': 'Software & Digital Tools', 'engage': 'Programs & Engagements',
}

# ===== Presentation =====
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]

# ===== Helpers =====
def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp

def add_rounded(slide, x, y, w, h, fill, radius=0.06):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.adjustments[0] = radius
    shp.shadow.inherit = False
    return shp

def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=TEXT,
             font=FONT_BODY, align='left', anchor='top'):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}[align]
    if anchor == 'middle': tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif anchor == 'bottom': tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    else: tf.vertical_anchor = MSO_ANCHOR.TOP
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.name = font
    r.font.color.rgb = color
    return tb

def add_paragraphs(slide, x, y, w, h, items, *, size=12, color=TEXT, font=FONT_BODY,
                   spacing=2, bold_first=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.space_after = Pt(spacing)
        r = p.add_run()
        r.text = item
        r.font.size = Pt(size); r.font.name = font
        r.font.color.rgb = color
        if bold_first and i == 0: r.font.bold = True
    return tb

def add_image_safely(slide, path, x, y, w, h):
    if path and os.path.exists(path):
        try:
            return slide.shapes.add_picture(path, x, y, width=w, height=h)
        except Exception as e:
            print(f'  IMG ERR {path}: {e}')
    return None

def footer(slide, page_num, total, section_label):
    add_rect(slide, 0, H - Inches(0.35), W, Inches(0.35), NAVY_DARK)
    add_text(slide, Inches(0.4), H - Inches(0.32), Inches(6), Inches(0.3),
             f'GIMI Product Catalogue 2026  |  {section_label}',
             size=9, color=RGBColor(0xCC, 0xCC, 0xCC), anchor='middle')
    add_text(slide, W - Inches(1.2), H - Inches(0.32), Inches(0.8), Inches(0.3),
             f'{page_num} / {total}',
             size=9, color=RGBColor(0xCC, 0xCC, 0xCC), align='right', anchor='middle')

# ===== Slide 1: Cover =====
def slide_cover():
    s = prs.slides.add_slide(BLANK)
    # background gradient — use solid navy with teal accent strip
    add_rect(s, 0, 0, W, H, NAVY)
    # Accent strip at left
    add_rect(s, 0, 0, Inches(0.4), H, TEAL)
    # GIMI logo top-right
    add_image_safely(s, LOGO_WHITE, W - Inches(2.6), Inches(0.5), Inches(2.0), Inches(0.66))
    # Big title
    add_text(s, Inches(1.2), Inches(2.4), Inches(11), Inches(1.2),
             'GIMI Product Catalogue', size=54, bold=True, color=WHITE, font=FONT_TITLE)
    add_text(s, Inches(1.2), Inches(3.7), Inches(11), Inches(0.6),
             '65 products to grow innovation capability at every level',
             size=22, color=RGBColor(0xCA, 0xDC, 0xFC), font=FONT_BODY)
    # Three pillar mini-tiles
    pillars = [('31', 'Develop people', TEAL),
               ('14', 'Build organisations', YELLOW),
               ('20', 'Equip and sustain', PINK)]
    for i, (num, lbl, col) in enumerate(pillars):
        x = Inches(1.2 + i*3.6)
        add_rect(s, x, Inches(4.8), Inches(3.3), Inches(1.4), col)
        add_text(s, x, Inches(4.85), Inches(3.3), Inches(0.7), num,
                 size=40, bold=True, color=NAVY, align='center')
        add_text(s, x, Inches(5.5), Inches(3.3), Inches(0.5), lbl,
                 size=14, color=NAVY, align='center', anchor='middle')
    add_text(s, Inches(1.2), Inches(6.6), Inches(11), Inches(0.4),
             'Last updated: May 2026   |   giminstitute.org',
             size=12, color=RGBColor(0xCC, 0xCC, 0xCC))

# ===== Slide 2: Minto overview =====
def slide_minto():
    s = prs.slides.add_slide(BLANK)
    # White bg
    add_rect(s, 0, 0, W, H, WHITE)
    # Navy header strip
    add_rect(s, 0, 0, W, Inches(1.4), NAVY)
    add_image_safely(s, LOGO_WHITE, Inches(0.5), Inches(0.35), Inches(1.6), Inches(0.55))
    add_text(s, Inches(2.4), Inches(0.4), Inches(10), Inches(0.5),
             'Catalogue at a glance', size=24, bold=True, color=WHITE)
    add_text(s, Inches(2.4), Inches(0.85), Inches(10), Inches(0.4),
             'How GIMI organises its 65 products', size=14, color=RGBColor(0xCA, 0xDC, 0xFC))
    # Headline / Minto answer
    add_text(s, Inches(0.6), Inches(1.7), Inches(12), Inches(0.7),
             'GIMI offers 65 products across 8 categories, grouped into 3 strategic pillars',
             size=22, bold=True, color=NAVY)

    # Three pillars
    pillars = [
        ('Develop people', TEAL, 31,
         ['Innovation Core (10)', 'Innovation Elective (10)', 'Future Foresight (5)',
          'Leadership Core (2)',  'Consulting Core (4)'],
         'Certifications across all career stages, from high-school students to chief innovation officers.'),
        ('Build the organisation', ORANGE, 14,
         ['Programs & Tools (14)'],
         'Assessments, software, programs, and engagements that drive innovation at the organisational level.'),
        ('Equip and sustain', PINK, 20,
         ['Books & Guides (13)', 'Memberships (7)'],
         'Reference materials and ongoing community access for individuals, companies, and academia.'),
    ]
    base_y = Inches(2.6)
    col_w = Inches(4.1); gap = Inches(0.15)
    base_x = Inches(0.55)
    for i, (name, color, count, items, blurb) in enumerate(pillars):
        x = base_x + i * (col_w + gap)
        # Card
        add_rect(s, x, base_y, col_w, Inches(4.0), GRAY_LIGHT)
        # Header strip in color
        add_rect(s, x, base_y, col_w, Inches(0.6), color)
        add_text(s, x, base_y, col_w, Inches(0.6), name,
                 size=16, bold=True, color=WHITE, align='center', anchor='middle')
        # Big count
        add_text(s, x, base_y + Inches(0.7), col_w, Inches(0.9),
                 str(count), size=56, bold=True, color=NAVY, align='center')
        add_text(s, x, base_y + Inches(1.55), col_w, Inches(0.3),
                 'products', size=12, color=GRAY, align='center')
        # Sections list
        y_items = base_y + Inches(2.0)
        for it in items:
            add_text(s, x + Inches(0.3), y_items, col_w - Inches(0.6), Inches(0.3),
                     '• ' + it, size=12, color=TEXT)
            y_items += Inches(0.32)
        # Blurb at bottom
        add_text(s, x + Inches(0.2), base_y + Inches(3.4), col_w - Inches(0.4), Inches(0.6),
                 blurb, size=10, color=GRAY)
    footer(s, 2, 62, 'Overview')

# ===== Section divider =====
def slide_section_divider(num, name, sub, count):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, W, H, NAVY)
    add_rect(s, 0, 0, Inches(0.4), H, TEAL)
    # Section number
    add_text(s, Inches(1.0), Inches(2.0), Inches(3.5), Inches(1.5),
             f'{num:02d}', size=140, bold=True, color=TEAL, font=FONT_TITLE)
    # Name
    add_text(s, Inches(4.6), Inches(2.6), Inches(8.0), Inches(1.0),
             name, size=44, bold=True, color=WHITE, font=FONT_TITLE)
    add_text(s, Inches(4.6), Inches(3.6), Inches(8.0), Inches(0.7),
             sub, size=18, color=RGBColor(0xCA, 0xDC, 0xFC))
    # Count
    add_text(s, Inches(4.6), Inches(4.5), Inches(8.0), Inches(0.6),
             f'{count} products in this section',
             size=14, color=YELLOW, bold=True)
    add_image_safely(s, LOGO_WHITE, W - Inches(2.6), Inches(0.5), Inches(2.0), Inches(0.66))

# ===== Product slide =====
def asset_path(badge_ref):
    if not badge_ref: return None
    p = os.path.join(r'C:\Users\aelha\Downloads', badge_ref.replace('/', os.sep))
    return p if os.path.exists(p) else None

def slide_product(p, page_num, total, section_label):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, W, H, WHITE)
    # Top bar (teal)
    add_rect(s, 0, 0, W, Inches(0.5), TEAL_VDARK)
    # Breadcrumb in top bar
    crumb = f'GIMI / {section_label}'
    if p.get('cluster') and p['cluster'] in CLUSTERS:
        crumb += f' / {CLUSTERS[p["cluster"]]}'
    if p.get('code'): crumb += f' / {p["code"]}'
    add_text(s, Inches(0.5), 0, Inches(12), Inches(0.5),
             crumb, size=12, color=WHITE, anchor='middle')

    # Title area
    add_text(s, Inches(0.5), Inches(0.7), Inches(12.3), Inches(0.8),
             p['title'], size=28, bold=True, color=NAVY, font=FONT_TITLE)
    # Tagline
    if p.get('tagline'):
        add_text(s, Inches(0.5), Inches(1.45), Inches(12.3), Inches(0.5),
                 p['tagline'], size=14, color=GRAY)

    # Yellow divider line
    add_rect(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(0.04), YELLOW)

    # Left: Badge in white card
    badge_path = asset_path(p.get('badge'))
    add_rounded(s, Inches(0.5), Inches(2.2), Inches(3.0), Inches(3.0), GRAY_LIGHT, radius=0.04)
    if badge_path:
        add_image_safely(s, badge_path, Inches(0.85), Inches(2.4), Inches(2.3), Inches(2.6))
    else:
        add_text(s, Inches(0.5), Inches(2.2), Inches(3.0), Inches(3.0),
                 p.get('code') or p['title'], size=20, bold=True, color=GRAY,
                 align='center', anchor='middle')

    # Right top: Description
    desc_x = Inches(3.8); desc_w = Inches(9.0)
    add_text(s, desc_x, Inches(2.2), desc_w, Inches(0.35),
             'DESCRIPTION', size=10, bold=True, color=TEAL_DARK)
    add_text(s, desc_x, Inches(2.55), desc_w, Inches(2.2),
             p.get('description') or p.get('tagline') or '', size=12, color=TEXT)

    # Bottom-left grid of metadata (4 cells)
    meta_y = Inches(5.4)
    cell_w = Inches(3.0); cell_h = Inches(0.9); gap = Inches(0.1)
    pairs = [
        ('Fee',      p.get('priceDisplay') or '—'),
        ('Duration', p.get('duration') or '—'),
        ('Modality', (p.get('modality') or '—').split('(')[0].strip() or '—'),
        ('For',      p.get('target') or '—'),
    ]
    for i, (k, v) in enumerate(pairs):
        x = Inches(0.5 + i*(3.0 + 0.1))
        add_rect(s, x, meta_y, cell_w, cell_h, GRAY_LIGHT)
        add_text(s, x + Inches(0.15), meta_y + Inches(0.08), cell_w, Inches(0.3),
                 k.upper(), size=9, bold=True, color=TEAL_DARK)
        add_text(s, x + Inches(0.15), meta_y + Inches(0.38), cell_w - Inches(0.2), Inches(0.5),
                 v, size=13, bold=True, color=NAVY)

    # Bottom bar
    add_rect(s, 0, H - Inches(0.35), W, Inches(0.35), NAVY_DARK)
    add_text(s, Inches(0.4), H - Inches(0.32), Inches(8), Inches(0.3),
             f'GIMI Product Catalogue 2026  |  {section_label}',
             size=9, color=RGBColor(0xCC,0xCC,0xCC), anchor='middle')
    add_text(s, W - Inches(1.2), H - Inches(0.32), Inches(0.8), Inches(0.3),
             f'{page_num} / {total}', size=9,
             color=RGBColor(0xCC,0xCC,0xCC), align='right', anchor='middle')

# ===== Memberships consolidated slide =====
def slide_memberships(page, total):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, W, H, WHITE)
    add_rect(s, 0, 0, W, Inches(0.5), TEAL_VDARK)
    add_text(s, Inches(0.5), 0, Inches(12), Inches(0.5),
             'GIMI / Memberships / All Tiers', size=12, color=WHITE, anchor='middle')
    add_text(s, Inches(0.5), Inches(0.7), Inches(12.3), Inches(0.7),
             'GIMI Memberships', size=28, bold=True, color=NAVY)
    add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.5),
             'Annual access tiers for individuals, companies, and academia',
             size=14, color=GRAY)
    add_rect(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(0.04), YELLOW)

    mems = [p for p in products if p['track'] == 'memb']
    cols = 4
    cw = Inches(3.0); ch = Inches(2.4)
    base_x = Inches(0.6); base_y = Inches(2.3); gx = Inches(0.18); gy = Inches(0.3)
    for i, m in enumerate(mems):
        row = i // cols; col = i % cols
        x = base_x + col * (cw + gx); y = base_y + row * (ch + gy)
        add_rect(s, x, y, cw, ch, GRAY_LIGHT)
        # badge
        bp = asset_path(m.get('badge'))
        if bp: add_image_safely(s, bp, x + Inches(0.85), y + Inches(0.15), Inches(1.3), Inches(1.3))
        # title
        title_short = m['title'].replace('Membership','M.').replace('Organizational','Org')
        add_text(s, x + Inches(0.15), y + Inches(1.5), cw - Inches(0.3), Inches(0.4),
                 m['title'], size=11, bold=True, color=NAVY, align='center')
        add_text(s, x + Inches(0.15), y + Inches(1.95), cw - Inches(0.3), Inches(0.35),
                 m.get('priceDisplay','—'), size=14, bold=True, color=TEAL_DARK, align='center')
    footer(s, page, total, 'Memberships')

# ===== Books consolidated slides per sub-cluster =====
def slide_book_cluster(cluster_id, label, page, total):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, W, H, WHITE)
    add_rect(s, 0, 0, W, Inches(0.5), TEAL_VDARK)
    add_text(s, Inches(0.5), 0, Inches(12), Inches(0.5),
             f'GIMI / Books & Guides / {label}', size=12, color=WHITE, anchor='middle')
    add_text(s, Inches(0.5), Inches(0.7), Inches(12.3), Inches(0.7),
             label, size=28, bold=True, color=NAVY)
    add_rect(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.04), YELLOW)

    items = [p for p in products if p['track'] == 'books' and p.get('cluster') == cluster_id]
    if not items: return
    cols = min(len(items), 4)
    cw = Inches(11.5 / cols); ch = Inches(4.5)
    base_x = Inches(0.5); base_y = Inches(1.9); gx = Inches(0.2)
    for i, b in enumerate(items):
        x = base_x + i * (cw + gx) if i < cols else base_x + (i - cols) * (cw + gx)
        y = base_y if i < cols else base_y + ch + Inches(0.2)
        add_rect(s, x, y, cw, ch, GRAY_LIGHT)
        bp = asset_path(b.get('badge'))
        if bp:
            add_image_safely(s, bp, x + Inches(0.15), y + Inches(0.15), cw - Inches(0.3), Inches(3.0))
        add_text(s, x + Inches(0.15), y + Inches(3.25), cw - Inches(0.3), Inches(0.7),
                 b['title'], size=11, bold=True, color=NAVY, align='center')
        add_text(s, x + Inches(0.15), y + Inches(4.0), cw - Inches(0.3), Inches(0.35),
                 b.get('priceDisplay','TBD'), size=12, bold=True, color=TEAL_DARK, align='center')
    footer(s, page, total, 'Books & Guides')

# ===== Closing =====
def slide_closing(page, total):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, W, H, NAVY)
    add_rect(s, 0, 0, Inches(0.4), H, TEAL)
    add_image_safely(s, LOGO_WHITE, W - Inches(2.6), Inches(0.5), Inches(2.0), Inches(0.66))
    add_text(s, Inches(1.0), Inches(2.2), Inches(11), Inches(1.2),
             'Thank you', size=72, bold=True, color=WHITE, font=FONT_TITLE)
    add_text(s, Inches(1.0), Inches(3.6), Inches(11), Inches(0.6),
             'Explore the full catalogue at giminstitute.org/store',
             size=22, color=YELLOW)
    add_text(s, Inches(1.0), Inches(4.6), Inches(11), Inches(0.5),
             'Global Innovation Management Institute',
             size=16, color=WHITE, bold=True)
    add_text(s, Inches(1.0), Inches(5.1), Inches(11), Inches(0.4),
             '110 Cambridge Street, Cambridge MA 02141',
             size=12, color=RGBColor(0xCC,0xCC,0xCC))
    add_text(s, Inches(1.0), Inches(5.45), Inches(11), Inches(0.4),
             'info@giminstitute.org',
             size=12, color=RGBColor(0xCC,0xCC,0xCC))

# ===== BUILD SEQUENCE =====
# Slide 1: Cover
slide_cover()
# Slide 2: Minto
slide_minto()

# Sections 1-8 with their products
total_slides = 62
page = 3  # 1=cover, 2=minto, 3+ start

section_order = [
    'icore', 'ielec', 'ff', 'lead', 'cons', 'books', 'tools', 'memb'
]

section_meta = {tid: (lbl, sub) for tid, lbl, sub in TRACKS}

for sec_num, sec in enumerate(section_order, start=1):
    section_label, sub = section_meta[sec]
    sec_products = [p for p in products if p['track'] == sec]

    # Section divider
    slide_section_divider(sec_num, section_label, sub, len(sec_products))
    page += 1

    if sec == 'memb':
        # Consolidated single slide
        slide_memberships(page, total_slides)
        page += 1
    elif sec == 'books':
        # 6 sub-clusters as separate slides
        for cluster_id, label in [
            ('imbok', 'IMBOK'),
            ('ffbok', 'Future Foresight Body of Knowledge'),
            ('iso-bk', 'ISO Standards'),
            ('stories', 'Innovation Stories'),
            ('tool-bk', 'Practitioner Tools'),
            ('journal', 'Journal'),
        ]:
            slide_book_cluster(cluster_id, label, page, total_slides)
            page += 1
    elif sec == 'tools':
        # Awards consolidated: Org + Individual on same slide
        # Build special: filter products
        awards = [p for p in sec_products if p['id'] in ('awards-individual','awards-org')]
        other_tools = [p for p in sec_products if p['id'] not in ('awards-individual','awards-org')]
        # Define logical sub-order
        order_ids = ['ipa-individual','ipa-corporate','mindset-index','oia-self','full-oia',
                     'akaio','eureka','service-tool','quickgrowth',
                     'academy','coaching','speech']
        for pid in order_ids:
            p = next((x for x in other_tools if x['id']==pid), None)
            if not p: continue
            slide_product(p, page, total_slides, section_label)
            page += 1
        # Awards combined slide
        if awards:
            # Use a combined custom slide
            s = prs.slides.add_slide(BLANK)
            add_rect(s, 0, 0, W, H, WHITE)
            add_rect(s, 0, 0, W, Inches(0.5), TEAL_VDARK)
            add_text(s, Inches(0.5), 0, Inches(12), Inches(0.5),
                     'GIMI / Programs & Tools / Recognition', size=12, color=WHITE, anchor='middle')
            add_text(s, Inches(0.5), Inches(0.7), Inches(12.3), Inches(0.7),
                     'GIMI Innovation Awards', size=28, bold=True, color=NAVY)
            add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.5),
                     'International recognition for groundbreaking innovation initiatives',
                     size=14, color=GRAY)
            add_rect(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(0.04), YELLOW)
            cw = Inches(6.0); ch = Inches(4.5)
            for i, a in enumerate(awards):
                x = Inches(0.7 + i*6.4); y = Inches(2.3)
                add_rect(s, x, y, cw, ch, GRAY_LIGHT)
                bp = asset_path(a.get('badge'))
                if bp:
                    add_image_safely(s, bp, x + Inches(2.0), y + Inches(0.3), Inches(2.0), Inches(2.0))
                add_text(s, x + Inches(0.3), y + Inches(2.5), cw - Inches(0.6), Inches(0.5),
                         a['title'], size=14, bold=True, color=NAVY, align='center')
                add_text(s, x + Inches(0.3), y + Inches(3.0), cw - Inches(0.6), Inches(0.4),
                         a.get('tagline',''), size=11, color=GRAY, align='center')
                add_text(s, x + Inches(0.3), y + Inches(3.6), cw - Inches(0.6), Inches(0.5),
                         a.get('priceDisplay','TBD'), size=18, bold=True, color=TEAL_DARK, align='center')
            footer(s, page, total_slides, 'Programs & Tools')
            page += 1
    else:
        # Standard: one slide per product
        # Order: flagships first, then by cluster declared order if applicable
        if sec == 'icore':
            cluster_order = ['cip','cip-mc','ccio','ccio-mc','impact','icore-train']
        elif sec == 'ff':
            cluster_order = ['ff-series','ff-train']
        else:
            cluster_order = []
        if cluster_order:
            ordered = []
            for cl in cluster_order:
                ordered.extend([p for p in sec_products if p.get('cluster')==cl])
            ordered.extend([p for p in sec_products if not p.get('cluster')])
            sec_products = ordered
        for p in sec_products:
            slide_product(p, page, total_slides, section_label)
            page += 1

# Closing
slide_closing(page, total_slides)

# Save
out = r'C:\Users\aelha\Downloads\GIMI_Product_Catalog_V51.pptx'
prs.save(out)
print('\nBuilt', len(prs.slides), 'slides ->', out)
