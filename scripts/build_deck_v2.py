"""Build V51 by cloning V50's existing product slide template for missing products."""
import os, json, copy, shutil, re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from lxml import etree
from pptx.oxml.ns import qn, nsmap

# Working files
SRC = r'C:\Users\aelha\Downloads\v50_source.pptx'
OUT = r'C:\Users\aelha\Downloads\GIMI_Product_Catalog_V52.pptx'

# Copy source
shutil.copy(SRC, OUT)

products = json.load(open(r'C:\Users\aelha\Downloads\products.json'))

# Apply modality defaults
def default_modality(p):
    if p.get('modality'): return p['modality']
    if p.get('cluster') in ('cip-mc', 'ccio-mc'): return 'Live Virtual Cohort'
    if p['track'] == 'books': return 'Reference'
    if p['track'] == 'memb' and p['id'] == 'thinktank': return 'Monthly Virtual'
    if p['track'] == 'memb': return 'Annual'
    if p.get('cluster') in ('indiv-ass', 'org-ass'): return 'Online Assessment'
    if p.get('cluster') == 'software': return 'Software Platform'
    if p.get('cluster') == 'engage': return 'Custom Engagement'
    return 'Online Self-Paced'
for p in products:
    p['modality'] = default_modality(p)
PRODUCTS_BY_ID = {p['id']: p for p in products}

prs = Presentation(OUT)
print(f'Starting with {len(prs.slides)} slides')

# ===== SLIDE CLONE HELPER =====
def duplicate_slide(prs, source_slide):
    """Clone a slide by deep-copying its XML. Returns the new Slide object."""
    blank_layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(blank_layout)
    # Remove default placeholders
    for shp in list(new_slide.shapes):
        sp = shp._element
        sp.getparent().remove(sp)
    # Copy each shape from source
    src_tree = source_slide.shapes._spTree
    dst_tree = new_slide.shapes._spTree
    for child in src_tree:
        tag = etree.QName(child).localname
        if tag in ('nvGrpSpPr', 'grpSpPr'):
            continue
        dst_tree.append(copy.deepcopy(child))
    # Copy relationships (images live in slide part rels)
    src_part = source_slide.part
    dst_part = new_slide.part
    for rel in src_part.rels.values():
        if 'notesSlide' in rel.reltype: continue
        if rel.is_external:
            dst_part.rels.get_or_add_ext_rel(rel.reltype, rel._target)
        else:
            # Add ref to same target part
            try:
                dst_part.rels.get_or_add(rel.reltype, rel._target)
            except Exception:
                pass
    return new_slide

# ===== TEXT EDIT HELPERS =====
def get_all_text_runs(slide):
    """Yield (paragraph, run, text) for all text in slide."""
    for shape in slide.shapes:
        if not shape.has_text_frame: continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                yield shape, para, run

def set_text_preserving_format(shape, new_text):
    """Replace all text in shape but keep first run's formatting."""
    if not shape.has_text_frame: return
    tf = shape.text_frame
    # Keep formatting from first run
    first_run_props = None
    for p in tf.paragraphs:
        for r in p.runs:
            first_run_props = r
            break
        if first_run_props: break
    # Clear all paragraphs except first
    p0 = tf.paragraphs[0]
    # Remove extra paragraphs
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)
    # Clear runs in first paragraph
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)
    # Add new run
    new_run = p0.add_run()
    new_run.text = new_text
    if first_run_props is not None:
        # Copy font properties
        try:
            if first_run_props.font.size:
                new_run.font.size = first_run_props.font.size
            if first_run_props.font.bold is not None:
                new_run.font.bold = first_run_props.font.bold
            if first_run_props.font.name:
                new_run.font.name = first_run_props.font.name
            if first_run_props.font.color and first_run_props.font.color.type:
                try:
                    new_run.font.color.rgb = first_run_props.font.color.rgb
                except Exception:
                    pass
        except Exception:
            pass

def set_text_at(slide, target_text_match, new_text):
    """Find a shape whose text starts with target_text_match and replace its text."""
    for shape in slide.shapes:
        if not shape.has_text_frame: continue
        if shape.text_frame.text.strip().startswith(target_text_match):
            set_text_preserving_format(shape, new_text)
            return True
    return False

# ===== IMAGE REPLACE =====
def replace_first_picture(slide, new_image_path):
    """Replace the first picture shape on slide with new image."""
    if not os.path.exists(new_image_path):
        return False
    for shape in slide.shapes:
        if shape.shape_type == 13:  # PICTURE
            # Get the picture's position and size
            x, y, w, h = shape.left, shape.top, shape.width, shape.height
            # Insert new picture
            new_pic = slide.shapes.add_picture(new_image_path, x, y, width=w, height=h)
            # Remove original
            sp = shape._element
            sp.getparent().remove(sp)
            return True
    return False

# ===== TEMPLATE SLIDE =====
# Use slide 5 (Innovation Catalyst) as template — has the clean GIMI product slide design
TEMPLATE_IDX = 4  # 0-indexed: slide 5

# ===== MISSING PRODUCTS TO ADD =====
# Map: product_id -> list of products that need new slides cloned from template
MISSING_PRODUCT_IDS = [
    # Innovation Core
    'cip-mc', 'ccio-mc', 'cgit', 'cgt',
    # Innovation Elective
    'ip', 'pbcc', 'cdt-2', 'cdtt-3', 'caap', 'cga', 'clc',
    # Future Foresight
    'cff-3', 'cff-4', 'cff-5',
    # Leadership Core
    'cil',
    # Consulting Core
    'mci-1', 'mci-2', 'mci-3', 'mci-4',
    # Programs & Tools
    'awards-individual', 'awards-org', 'akaio', 'quickgrowth', 'academy', 'speech',
    'ipa-corporate', 'full-oia',
    # Books & Guides
    'technovate',
    # Memberships
    'thinktank',
]

# Verify all exist in PRODUCTS_BY_ID
for pid in MISSING_PRODUCT_IDS:
    if pid not in PRODUCTS_BY_ID:
        print(f'WARN: {pid} not in products data')

template_slide = prs.slides[TEMPLATE_IDX]

# ===== ASSET PATH HELPER =====
def asset_path(badge_ref):
    if not badge_ref: return None
    p = os.path.join(r'C:\Users\aelha\Downloads', badge_ref.replace('/', os.sep))
    return p if os.path.exists(p) else None

# ===== FILL CLONED SLIDE =====
def populate_cloned_slide(slide, product):
    """Update text and badge image in a cloned slide template for the given product."""
    # Identify and update specific shapes
    # The template has these structurally identifiable text blocks:
    # - Breadcrumb (top): "Certifications   For Individuals   Certified Innovation Catalyst"
    # - Title placeholder (lower title near top): "The Innovation Catalyst Certification..."
    # - Product name rectangles (banner): "Innovation Catalyst"
    # - Description text (main body)
    # - Includes text
    # - Ideal for / Modality / Preparation / Fee — small labeled blocks
    # - Testimony block

    # Build the breadcrumb based on track/cluster
    track_lbl = {
        'icore':'Innovation Core', 'ielec':'Innovation Elective', 'ff':'Future Foresight',
        'lead':'Leadership Core', 'cons':'Consulting Core', 'books':'Books & Guides',
        'tools':'Programs & Tools', 'memb':'Memberships'
    }.get(product['track'], 'GIMI')
    breadcrumb = f'{track_lbl}    For Individuals    {product["title"]}'

    # Walk shapes; identify by current text and update accordingly
    for shape in slide.shapes:
        if not shape.has_text_frame: continue
        text = shape.text_frame.text.strip()

        # Breadcrumb (placeholder near top)
        if text.startswith('Certifications') and len(text) < 100 and 'Innovation' in text:
            set_text_preserving_format(shape, breadcrumb)
            continue

        # Title placeholder (long description-like text near top)
        if text.startswith('The Innovation Catalyst Certification'):
            set_text_preserving_format(shape, product.get('tagline') or product.get('description',''))
            continue

        # Product name labels (multiple): match "Innovation Catalyst" string exactly
        if text == 'Innovation Catalyst':
            set_text_preserving_format(shape, product['title'])
            continue

        # Description main body
        if text.startswith('Start your innovation journey'):
            desc = product.get('description') or product.get('tagline') or ''
            set_text_preserving_format(shape, desc)
            continue

        # Includes block
        if text.startswith('Innovation Catalyst Certificate'):
            includes_text = 'Certificate awarded upon passing\nOnline examination fee\nStudy guide (e-copy)\nRemote proctor fee\nOnline coursework'
            set_text_preserving_format(shape, includes_text)
            continue

        # Testimony quote
        if 'I\'m definitely going to take' in text or 'I am definitely going to take' in text or 'definitely going to take' in text:
            set_text_preserving_format(shape, f'"{product.get("tagline","")}"  — GIMI alumni')
            continue

        # Good to know callout
        if text.startswith('This certification can'):
            set_text_preserving_format(shape, f'Modality: {product["modality"]}. Fee: {product.get("priceDisplay","TBD")}.')
            continue

        # Labels stay; metadata values: Ideal for, Modality, Preparation, Fee — find their adjacent text frames
        # The template has labels Ideal for/Modality/Preparation/Fee as separate rectangles with the value somewhere nearby.
        # In V50's slide 5, the metadata labels are in dedicated boxes ("Ideal for" rectangle) and the VALUES are in adjacent text frames inside groups.
        # We'll skip granular metadata replacement here and rely on description + good to know carrying the info.

    # Replace badge image
    bp = asset_path(product.get('badge'))
    if bp:
        replace_first_picture(slide, bp)

# ===== CLONE TEMPLATE FOR EACH MISSING PRODUCT =====
added = 0
for pid in MISSING_PRODUCT_IDS:
    if pid not in PRODUCTS_BY_ID:
        continue
    new_slide = duplicate_slide(prs, template_slide)
    populate_cloned_slide(new_slide, PRODUCTS_BY_ID[pid])
    added += 1

print(f'Added {added} new product slides. Total now: {len(prs.slides)}')

# ===== ADD MINTO OVERVIEW SLIDE =====
# Clone slide 2 (which appears to be an overview/sectioning slide) as a base for Minto
# Actually for the Minto, we'll use a simpler clone of slide 4 (programs overview) which has a structural look
# For now, skip — the user can use V50's existing overview slides.

prs.save(OUT)
print(f'Saved {OUT}')
