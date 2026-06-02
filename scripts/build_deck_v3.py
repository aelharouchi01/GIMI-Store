"""GIMI Product Catalogue V53 — clean rebuild
- V50-style product slides with full info (description + includes + 4 metadata cells)
- Books consolidated to 3 slides, Programs to 4, Memberships to 1
- Aspect-ratio preserving badges
- All text embedded inside shapes (move-as-one)
- Explicit white space between every box (no flush edges)
"""
import os, json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

# ============ BRAND ============
TEAL       = RGBColor(0x2A, 0xAC, 0xB5)
TEAL_DARK  = RGBColor(0x1E, 0x8A, 0x93)
TEAL_VDARK = RGBColor(0x10, 0x4E, 0x60)
YELLOW     = RGBColor(0xC9, 0xD9, 0x40)
NAVY       = RGBColor(0x1B, 0x1F, 0x3B)
NAVY_DARK  = RGBColor(0x13, 0x16, 0x2C)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY       = RGBColor(0x6B, 0x72, 0x80)
GRAY_LIGHT = RGBColor(0xF0, 0xF4, 0xF6)
GRAY_BG    = RGBColor(0xFA, 0xFB, 0xFC)
TEXT       = RGBColor(0x1F, 0x29, 0x37)
PINK       = RGBColor(0xE1, 0x5A, 0x78)
ORANGE     = RGBColor(0xF0, 0x87, 0x1E)
FONT = 'Calibri'

LOGO_WHITE = r'C:\Users\aelha\Downloads\gimi-store-assets\gimi-logo-white.png'

# ============ DATA ============
products = json.load(open(r'C:\Users\aelha\Downloads\products.json'))
def default_modality(p):
    if p.get('modality'): return p['modality']
    if p.get('cluster') in ('cip-mc','ccio-mc'): return 'Live Virtual Cohort'
    if p['track']=='books': return 'Reference'
    if p['track']=='memb' and p['id']=='thinktank': return 'Monthly Virtual'
    if p['track']=='memb': return 'Annual'
    if p.get('cluster') in ('indiv-ass','org-ass'): return 'Online Self-Assessment'
    if p.get('cluster')=='software': return 'Software Platform'
    if p.get('cluster')=='engage': return 'Custom Engagement'
    return 'Online Self-Paced'
for p in products: p['modality'] = default_modality(p)
PB = {p['id']: p for p in products}

TRACKS = [
    ('icore', 'Innovation Core',     'The flagship innovation management pathway', TEAL),
    ('ielec', 'Innovation Elective', 'Specialist tracks across disciplines',        ORANGE),
    ('ff',    'Future Foresight',    'Foresight and trainer certifications',        PINK),
    ('lead',  'Leadership Core',     'Leadership-oriented certifications',           YELLOW),
    ('cons',  'Consulting Core',     'Management Consulting Institute pathway',     NAVY),
    ('books', 'Books & Guides',      'Published books and reference materials',     GRAY),
    ('tools', 'Programs & Tools',    'Assessments, software, programs',             TEAL_DARK),
    ('memb',  'Memberships',         'Annual access tiers',                          NAVY_DARK),
]
CLUSTERS = {
    'cip':'CIP series','cip-mc':'CIP-Masterclass','ccio':'CCIO series','ccio-mc':'CCIO-Masterclass',
    'impact':'GIMI Impact','icore-train':'Train the Trainer',
    'ff-series':'Foresight series','ff-train':'Train the Trainer',
    'imbok':'IMBOK','ffbok':'Future Foresight BOK','iso-bk':'ISO Standards',
    'stories':'Innovation Stories','tool-bk':'Practitioner Tools','journal':'Journal',
    'indiv-ass':'Individual Assessments','org-ass':'Organisational Assessments',
    'software':'Software & Digital Tools','engage':'Programs & Engagements',
}

# ============ SETUP PRESENTATION ============
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

# ============ HELPERS ============
def add_rect(slide, x, y, w, h, fill, *, line=None, text=None, size=14, bold=False,
             color=TEXT, align='left', anchor='top', rounded=False, radius=0.06,
             margin=0.05, italic=False):
    """Add a rectangle (or rounded rectangle) with optional embedded text."""
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(0.5)
    if rounded:
        shp.adjustments[0] = radius
    try: shp.shadow.inherit = False
    except Exception: pass
    if text is not None:
        tf = shp.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(margin); tf.margin_right = Inches(margin)
        tf.margin_top = Inches(margin); tf.margin_bottom = Inches(margin)
        anchors = {'top': MSO_ANCHOR.TOP, 'middle': MSO_ANCHOR.MIDDLE, 'bottom': MSO_ANCHOR.BOTTOM}
        tf.vertical_anchor = anchors[anchor]
        p = tf.paragraphs[0]
        p.alignment = {'left':PP_ALIGN.LEFT,'center':PP_ALIGN.CENTER,'right':PP_ALIGN.RIGHT}[align]
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.name = FONT; r.font.color.rgb = color
    return shp

def add_rect_lines(slide, x, y, w, h, fill, lines, *, size=11, color=TEXT,
                   bullet=False, anchor='top', line_spacing=4, margin=0.1):
    """Rectangle with multi-line embedded text. Bullets get proper hanging indent."""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    try: shp.shadow.inherit = False
    except Exception: pass
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(margin); tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin); tf.margin_bottom = Inches(margin)
    anchors = {'top': MSO_ANCHOR.TOP, 'middle': MSO_ANCHOR.MIDDLE, 'bottom': MSO_ANCHOR.BOTTOM}
    tf.vertical_anchor = anchors[anchor]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(line_spacing)
        if bullet:
            # Hanging indent so wrapped lines align with first character (not bullet)
            from pptx.oxml.ns import qn
            pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
            pPr.set('marL', '228600')   # 0.25" left margin
            pPr.set('indent', '-228600') # -0.25" first-line indent
        r = p.add_run()
        r.text = ('• ' + line) if bullet else line
        r.font.size = Pt(size); r.font.name = FONT; r.font.color.rgb = color
    return shp

def add_picture_fitted(slide, image_path, x, y, max_w, max_h):
    """Add image inside box (x,y,max_w,max_h) preserving aspect ratio. Centered."""
    if not image_path or not os.path.exists(image_path): return None
    try:
        im = Image.open(image_path)
        iw, ih = im.size
    except Exception:
        iw, ih = 1, 1
    aspect = iw / ih if ih else 1
    box_aspect = max_w / max_h if max_h else 1
    if aspect > box_aspect:
        w = max_w; h = int(max_w / aspect)
    else:
        h = max_h; w = int(max_h * aspect)
    px = x + (max_w - w) // 2
    py = y + (max_h - h) // 2
    try:
        return slide.shapes.add_picture(image_path, px, py, width=w, height=h)
    except Exception as e:
        print(f'  IMG ERR {image_path}: {e}')
        return None

def asset_path(badge_ref):
    if not badge_ref: return None
    p = os.path.join(r'C:\Users\aelha\Downloads', badge_ref.replace('/', os.sep))
    return p if os.path.exists(p) else None

def footer(slide, page, total, section_label):
    add_rect(slide, 0, H - Inches(0.4), W, Inches(0.4), NAVY_DARK,
             text=f'GIMI Product Catalogue 2026   |   {section_label}',
             size=9, color=RGBColor(0xCC,0xCC,0xCC), align='left',
             anchor='middle', margin=0.4)
    add_rect(slide, W - Inches(1.3), H - Inches(0.4), Inches(1.0), Inches(0.4),
             NAVY_DARK, text=f'{page} / {total}',
             size=9, color=RGBColor(0xCC,0xCC,0xCC), align='right',
             anchor='middle', margin=0.15)

# ============ STANDARDIZED INCLUDES (per type) ============
def get_type(p):
    if p.get('cluster') in ('cip-mc','ccio-mc'): return 'masterclass'
    if p['track'] in ('icore','ielec','ff','lead','cons'): return 'cert'
    if p['track']=='books': return 'book'
    if p['track']=='memb': return 'community' if p['id']=='thinktank' else 'membership'
    if p.get('cluster') in ('indiv-ass','org-ass'): return 'assessment'
    if p.get('cluster')=='software': return 'software'
    if p.get('cluster')=='engage': return 'program'
    return 'product'

def standard_includes(p):
    t = get_type(p)
    if t=='cert': return [
        'Certificate awarded upon passing',
        'Online proctored examination fee (70% pass)',
        'Body of Knowledge / Study Guide (e-copy)',
        'Remote proctor fee + coursework',
        'Recertification: 3-year cycle, 60 CPE hrs',
    ]
    if t=='masterclass': return [
        'Live virtual cohort sessions (7 weeks)',
        'Final examination fee (70% pass)',
        'Official Study Guide + ISO Panel Webinar',
        'Certificate upon passing',
        'Recertification: 3-year cycle, 60 CPE hrs',
    ]
    if t=='assessment': return [
        'Self-assessment online survey',
        'Automated personalised report',
        'Strengths, gaps, and persona profile',
        'Actionable improvement roadmap',
    ]
    if t=='software': return [
        'Platform access for subscription period',
        'User onboarding and documentation',
        'Email support and product updates',
        'Periodic feature releases',
    ]
    if t=='program': return [
        'Tailored engagement scope at kickoff',
        'Structured methodology and timeline',
        'Documented outputs and deliverables',
        'GIMI expert facilitation',
    ]
    return ['Reference material', 'Digital delivery', 'Innovation case studies']

def good_to_know(p):
    t = get_type(p)
    if t=='cert': return f'Recertification: 3-year cycle, 60 CPE hrs. Pass threshold: 70%.'
    if t=='masterclass': return f'Includes ISO 56002/56003/56004 framework modules. Pass threshold: 70%.'
    if t=='assessment': return 'Outputs an automated PDF report with persona/maturity profile.'
    if t=='software': return 'Tiered access; contact GIMI for enterprise terms.'
    if t=='program': return 'Custom scope and timeline agreed at kickoff.'
    return ''

# ============ LEARNING OBJECTIVES / SKILLS / CAREER ROLES ============
# Hand-curated per cert (or per family) with sensible defaults derived from type.

PRODUCT_OUTCOMES = {
    # Innovation Core
    'cip-0':   {'objectives':['Understand business innovation fundamentals','Apply key idea-generation tools','Frame innovation challenges'],
                'skills':['Innovation Mindset','Idea Generation','Problem Framing','Creative Thinking'],
                'roles':['Innovation Champion','Junior Innovator','Project Contributor']},
    'cip-1':   {'objectives':['Master the GIMI ideation toolkit','Generate breakthrough ideas at scale','Apply ISO 56000-aligned frameworks'],
                'skills':['Ideation','Innovation Strategy','Concept Evaluation','Lean Experimentation'],
                'roles':['Innovation Associate','Product Innovator','Innovation Analyst']},
    'cip-2':   {'objectives':['Lead a real-world innovation case study','Document innovation outcomes','Demonstrate growth-idea generation'],
                'skills':['Innovation Project Management','Case Documentation','Growth Idea Generation','Stakeholder Engagement'],
                'roles':['Innovation Master','Innovation Lead','Senior Innovation Analyst']},
    'cip-mc':  {'objectives':['Master end-to-end innovation management','Apply ISO 56000 framework','Lead cross-functional teams'],
                'skills':['Innovation Management','ISO 56000 Frameworks','Cohort Leadership','Strategic Innovation'],
                'roles':['CIP Professional','Innovation Manager','Strategy Consultant']},
    'ccio-1':  {'objectives':['Design organisational innovation systems','Govern innovation strategy','Lead innovation teams'],
                'skills':['Innovation Governance','Portfolio Management','Strategic Foresight','Team Leadership'],
                'roles':['Innovation Manager','Innovation Director','Head of R&D']},
    'ccio-2':  {'objectives':['Build innovation ecosystems at scale','Manage innovation portfolios','Drive enterprise transformation'],
                'skills':['Ecosystem Design','Portfolio Strategy','Enterprise Leadership','Innovation Investment'],
                'roles':['Chief Innovation Officer','VP Innovation','Innovation Practice Lead']},
    'ccio-mc': {'objectives':['Lead enterprise innovation programs','Mentor senior leaders','Govern ISO-aligned systems'],
                'skills':['Executive Leadership','ISO Compliance','Mentorship','Strategic Communication'],
                'roles':['Chief Innovation Officer','SVP Innovation','Board Advisor']},
    'cgis-1':  {'objectives':['Apply creativity to global challenges','Develop a compassionate innovation mindset','Lead student innovation projects'],
                'skills':['Creative Problem Solving','Empathy','Project Initiative','Team Collaboration'],
                'roles':['Student Innovator','Project Leader','Future Entrepreneur']},
    'cgit':    {'objectives':['Teach innovation to high schoolers','Facilitate problem-solving','Build classroom innovation culture'],
                'skills':['Innovation Pedagogy','Facilitation','Mentorship','Curriculum Design'],
                'roles':['Innovation Teacher','School Coordinator','Education Mentor']},
    'cgt':     {'objectives':['Train others in GIMI methodologies','Deliver certified courses','Build a training business'],
                'skills':['Adult Education','GIMI Frameworks','Public Speaking','Course Delivery'],
                'roles':['GIMI Trainer','Innovation Coach','Independent Consultant']},

    # Innovation Elective
    'ip':      {'objectives':['Become a proactive driver of value','Connect knowledge and research','Build foundational innovation skills'],
                'skills':['Connection Thinking','Value Discovery','Self-direction','Innovation Awareness'],
                'roles':['Innovation Advocate','Curious Professional','Knowledge Connector']},
    'pbcc':    {'objectives':['Develop a structured problem-solving approach','Generate creative solutions','Pitch innovation concepts'],
                'skills':['5-Step Ideation','Creative Problem Solving','Concept Pitching','Divergent Thinking'],
                'roles':['Problem Solver','Innovation Catalyst','Student Founder']},
    'cdt-1':   {'objectives':['Master human-centered design','Build empathy with users','Prototype creative solutions'],
                'skills':['Design Thinking','User Empathy','Prototyping','Insight Generation'],
                'roles':['Design Thinker','UX Researcher','Innovation Designer']},
    'cdt-2':   {'objectives':['Apply Design Thinking at scale','Document innovation outcomes','Lead design challenges'],
                'skills':['Design Leadership','Workshop Facilitation','Case Documentation','Iterative Design'],
                'roles':['Senior Designer','Design Lead','Innovation Consultant']},
    'cdtt-3':  {'objectives':['Train others in Design Thinking','Facilitate design workshops','Build design culture'],
                'skills':['Design Facilitation','Train-the-Trainer','Workshop Design','Coaching'],
                'roles':['Design Thinking Trainer','Workshop Facilitator','Design Coach']},
    'caap':    {'objectives':['Operationalise AI in your team','Design AI agent solutions','Build executive-ready prototypes'],
                'skills':['AI Strategy','Agent Design','Innovation 5.0','Prototyping'],
                'roles':['AI Practitioner','Innovation Lead','Solution Architect']},
    'cga':     {'objectives':['Audit innovation systems','Apply ISO 56001 standard','Assess organisational maturity'],
                'skills':['ISO 56001 Auditing','Compliance Assessment','Maturity Evaluation','Reporting'],
                'roles':['GIMI Auditor','Innovation Auditor','Compliance Lead']},
    'cime':    {'objectives':['Implement ISO 56000-aligned systems','Design innovation management systems','Lead ISO transformations'],
                'skills':['ISO 56000 Mastery','Systems Design','Compliance Strategy','Implementation Leadership'],
                'roles':['ISO Expert','Innovation Systems Architect','Senior Consultant']},
    'ctc':     {'objectives':['Leverage 130+ breakthrough technologies','Identify emerging trends','Craft future-ready strategies'],
                'skills':['Technology Foresight','Trend Identification','Strategy Design','Tech Scouting'],
                'roles':['Technology Catalyst','Trends Analyst','Strategy Consultant']},
    'clc':     {'objectives':['Lead the shift to a longevity society','Understand brain & body health innovation','Champion longevity products'],
                'skills':['Longevity Innovation','Trend Analysis','Mindset Shift Leadership','Product Vision'],
                'roles':['Longevity Catalyst','HealthTech Strategist','Wellness Innovator']},

    # Future Foresight
    'cff-1':   {'objectives':['Question current assumptions','Apply foresight fundamentals','Explore schools of thought'],
                'skills':['Strategic Foresight','Trend Analysis','Critical Thinking','Scenario Awareness'],
                'roles':['Foresight Professional','Strategy Analyst','Research Lead']},
    'cff-2':   {'objectives':['Master the S4 Futures framework','Run end-to-end foresight projects','Apply multiple methodologies'],
                'skills':['Scoping','Scanning','Scenarios','Strategy'],
                'roles':['Foresight Leader','Senior Strategist','Foresight Manager']},
    'cff-3':   {'objectives':['Lead a Foresight Project','Integrate methodological phases','Drive strategic foresight'],
                'skills':['Foresight Project Leadership','Methodology Mastery','Executive Communication','Decision Support'],
                'roles':['Chief Foresight Officer','VP Strategy','Foresight Director']},
    'cff-4':   {'objectives':['Train others in foresight L1','Facilitate foresight workshops','Build foresight capability'],
                'skills':['Foresight Pedagogy','Train-the-Trainer','Facilitation','Coaching'],
                'roles':['Foresight Trainer','Workshop Lead','Foresight Coach']},
    'cff-5':   {'objectives':['Train advanced foresight practitioners','Deliver Level 2 certifications','Mentor foresight projects'],
                'skills':['Advanced Foresight','Senior Coaching','Mentorship','Project Review'],
                'roles':['Senior Foresight Trainer','Foresight Master Coach','Foresight Director']},

    # Leadership Core
    'clf-1':   {'objectives':['Lead collaboratively','Set bigger, bolder targets','Drive results through teams'],
                'skills':['Collaborative Leadership','Powers of AND','Vision Setting','Team Influence'],
                'roles':['Future Leader','Senior Manager','Director']},
    'cil':     {'objectives':['Build the SCD framework mindset','Lead repeatable innovation','Drive measurable change'],
                'skills':['Strategy-Capacity-Discipline','Innovation Leadership','Repeatable Process Design','Culture Building'],
                'roles':['Innovation Leader','Team Lead','Change Maker']},

    # Consulting Core
    'mci-1':   {'objectives':['Think like a consultant','Apply analytical techniques','Solve structured problems'],
                'skills':['Analytical Thinking','Problem Solving','Research','Structured Communication'],
                'roles':['Consulting Analyst','Junior Consultant','Business Analyst']},
    'mci-2':   {'objectives':['Lead a sponsored consulting project','Apply SCQ analysis','Deliver client outcomes'],
                'skills':['Project Delivery','SCQ Analysis','Client Engagement','Peer Feedback'],
                'roles':['Management Consultant','Engagement Lead','Senior Analyst']},
    'mci-3':   {'objectives':['Master Strategy Formulation','Apply REALIZE methodology','Bridge strategy to execution'],
                'skills':['Strategy Formulation','Execution Planning','REALIZE Methodology','Leadership'],
                'roles':['Consulting Manager','Strategy Lead','Engagement Manager']},
    'mci-4':   {'objectives':['Shape the consulting practice','Lead complex engagements','Drive industry impact'],
                'skills':['Visionary Leadership','Advanced Strategy','Industry Influence','Practice Building'],
                'roles':['Consulting Partner','Practice Lead','Industry Director']},
}

def get_outcomes(p):
    """Return (objectives, skills, roles) lists for a product. Falls back to type-based defaults."""
    if p['id'] in PRODUCT_OUTCOMES:
        d = PRODUCT_OUTCOMES[p['id']]
        return d['objectives'], d['skills'], d['roles']
    # Defaults by type for products not in the table
    t = get_type(p)
    if t == 'masterclass':
        return (['Master innovation management at scale','Apply ISO 56000 standards','Lead executive programs'],
                ['Strategic Innovation','ISO Mastery','Executive Communication','Cohort Leadership'],
                ['Senior Innovation Leader','Innovation Director','Practice Lead'])
    if t == 'assessment':
        return (['Understand your innovation profile','Identify strengths and gaps','Receive a tailored roadmap'],
                ['Self-awareness','Diagnostic Reading','Growth Planning'],
                ['Self-driven Learner','Career Planner','Innovation Champion'])
    if t == 'software':
        return (['Use the platform to accelerate innovation','Generate scenarios with AI','Match with the right partners'],
                ['Platform Adoption','AI Tooling','Ecosystem Building'],
                ['Innovation Lead','Strategy Owner','Platform Champion'])
    if t == 'program':
        return (['Design and run an engagement','Drive measurable outcomes','Build internal capability'],
                ['Program Design','Capability Building','Outcome Measurement'],
                ['Program Sponsor','Engagement Lead','Organisational Champion'])
    return (['Build foundational innovation knowledge','Apply best practices','Drive personal impact'],
            ['Innovation Awareness','Application','Communication'],
            ['Innovation Contributor','Team Member','Project Participant'])

# ============ COVER ============
def slide_cover():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, W, H, NAVY)
    add_rect(s, 0, 0, Inches(0.4), H, TEAL)
    s.shapes.add_picture(LOGO_WHITE, W - Inches(2.6), Inches(0.5),
                         width=Inches(2.0), height=Inches(0.66))
    add_rect(s, Inches(1.2), Inches(2.2), Inches(11.5), Inches(1.3), NAVY,
             text='GIMI Product Catalogue',
             size=54, bold=True, color=WHITE, align='left', anchor='middle')
    add_rect(s, Inches(1.2), Inches(3.55), Inches(11.5), Inches(0.6), NAVY,
             text='65 products to grow innovation capability at every level',
             size=22, color=RGBColor(0xCA,0xDC,0xFC), align='left', anchor='middle')
    # 3 pillar tiles with white space between
    pillars = [('31','Develop people', TEAL),
               ('14','Build organisations', YELLOW),
               ('20','Equip and sustain', PINK)]
    tile_w = Inches(3.4); gap = Inches(0.25); start_x = Inches(1.2)
    for i,(num,lbl,col) in enumerate(pillars):
        x = start_x + i*(tile_w + gap)
        add_rect(s, x, Inches(4.8), tile_w, Inches(1.4), col,
                 text=num, size=44, bold=True, color=NAVY,
                 align='center', anchor='middle')
        # Label below the tile (separate shape)
        add_rect(s, x, Inches(6.25), tile_w, Inches(0.4), NAVY,
                 text=lbl, size=14, bold=True, color=WHITE,
                 align='center', anchor='middle')
    add_rect(s, Inches(1.2), Inches(6.85), Inches(11.5), Inches(0.4), NAVY,
             text='Last updated: May 2026   |   giminstitute.org',
             size=12, color=RGBColor(0xCC,0xCC,0xCC), align='left', anchor='middle')

# ============ MINTO OVERVIEW ============
def slide_minto():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, W, H, WHITE)
    # Header strip
    add_rect(s, 0, 0, W, Inches(1.3), NAVY)
    s.shapes.add_picture(LOGO_WHITE, Inches(0.5), Inches(0.35),
                         width=Inches(1.6), height=Inches(0.53))
    add_rect(s, Inches(2.4), Inches(0.3), Inches(10), Inches(0.5), NAVY,
             text='Catalogue at a glance', size=22, bold=True, color=WHITE)
    add_rect(s, Inches(2.4), Inches(0.78), Inches(10), Inches(0.4), NAVY,
             text='How GIMI organises its 65 products',
             size=13, color=RGBColor(0xCA,0xDC,0xFC))
    # Minto headline
    add_rect(s, Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.7), WHITE,
             text='GIMI offers 65 products across 8 categories, grouped into 3 strategic pillars',
             size=20, bold=True, color=NAVY, align='left', anchor='top')

    # 3 pillar cards with white space between
    pillars = [
        ('Develop people', TEAL, 31,
         ['Innovation Core (10)', 'Innovation Elective (10)', 'Future Foresight (5)',
          'Leadership Core (2)', 'Consulting Core (4)'],
         'Certifications across all career stages, from high-school students to chief innovation officers.'),
        ('Build the organisation', ORANGE, 14,
         ['Programs & Tools (14)'],
         'Assessments, software, programs, and engagements driving innovation at the organisational level.'),
        ('Equip and sustain', PINK, 20,
         ['Books & Guides (13)', 'Memberships (7)'],
         'Reference materials and community access for individuals, companies, and academia.'),
    ]
    card_w = Inches(4.0); gap = Inches(0.2); base_x = Inches(0.6); base_y = Inches(2.5)
    card_h = Inches(4.5)
    for i,(name,color,count,items,blurb) in enumerate(pillars):
        x = base_x + i*(card_w + gap)
        # Header bar
        add_rect(s, x, base_y, card_w, Inches(0.6), color,
                 text=name, size=15, bold=True, color=WHITE, align='center', anchor='middle')
        # Body (separated by tiny gap)
        body_y = base_y + Inches(0.7)
        body_h = card_h - Inches(0.7)
        add_rect(s, x, body_y, card_w, body_h, GRAY_LIGHT)
        # Big count - separate shape inside body
        add_rect(s, x + Inches(0.2), body_y + Inches(0.2), card_w - Inches(0.4), Inches(1.0),
                 GRAY_LIGHT, text=str(count), size=54, bold=True, color=NAVY,
                 align='center', anchor='middle')
        add_rect(s, x + Inches(0.2), body_y + Inches(1.25), card_w - Inches(0.4), Inches(0.3),
                 GRAY_LIGHT, text='products', size=12, color=GRAY,
                 align='center', anchor='middle')
        # Section list as chip cards
        chip_y = body_y + Inches(1.65)
        chip_w = card_w - Inches(0.6); chip_h = Inches(0.28); chip_gap = Inches(0.06)
        for it in items:
            add_rect(s, x + Inches(0.3), chip_y, chip_w, chip_h, WHITE,
                     text=it, size=11, color=NAVY, align='center',
                     anchor='middle', margin=0.05)
            chip_y += chip_h + chip_gap
        # Blurb at bottom
        add_rect(s, x + Inches(0.3), body_y + body_h - Inches(0.85),
                 card_w - Inches(0.6), Inches(0.75), GRAY_LIGHT,
                 text=blurb, size=10, color=GRAY, anchor='top')
    footer(s, 2, 50, 'Overview')

# ============ SECTION DIVIDER ============
def slide_section_divider(num, name, sub, count, color, page, total):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, W, H, NAVY)
    add_rect(s, 0, 0, Inches(0.4), H, color)
    s.shapes.add_picture(LOGO_WHITE, W - Inches(2.6), Inches(0.5),
                         width=Inches(2.0), height=Inches(0.66))
    add_rect(s, Inches(1.0), Inches(1.8), Inches(3.5), Inches(2.4), NAVY,
             text=f'{num:02d}', size=110, bold=True, color=color,
             align='left', anchor='middle')
    add_rect(s, Inches(4.6), Inches(2.5), Inches(8.0), Inches(1.1), NAVY,
             text=name, size=42, bold=True, color=WHITE, align='left', anchor='middle')
    add_rect(s, Inches(4.6), Inches(3.7), Inches(8.0), Inches(0.6), NAVY,
             text=sub, size=17, color=RGBColor(0xCA,0xDC,0xFC), align='left')
    add_rect(s, Inches(4.6), Inches(4.5), Inches(8.0), Inches(0.5), NAVY,
             text=f'{count} products in this section',
             size=14, bold=True, color=YELLOW, align='left')

# ============ PRODUCT SLIDE ============
def slide_product(p, page, total, section_label, accent):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, W, H, WHITE)

    # ----- Top bar (breadcrumb) -----
    crumb = f'GIMI  /  {section_label}'
    if p.get('cluster') and p['cluster'] in CLUSTERS:
        crumb += f'  /  {CLUSTERS[p["cluster"]]}'
    if p.get('code'): crumb += f'  /  {p["code"]}'
    add_rect(s, 0, 0, W, Inches(0.4), TEAL_VDARK,
             text=crumb, size=11, color=WHITE, align='left', anchor='middle', margin=0.4)

    # ----- Title -----
    add_rect(s, Inches(0.5), Inches(0.5), Inches(12.3), Inches(0.55), WHITE,
             text=p['title'], size=22, bold=True, color=NAVY, anchor='middle')

    # ----- Storyline (black, not italic, bolder feel) -----
    storyline = p.get('tagline') or p.get('description','').split('.')[0]
    add_rect(s, Inches(0.5), Inches(1.07), Inches(12.3), Inches(0.38), WHITE,
             text=storyline, size=13, color=TEXT, italic=False, anchor='middle')

    # ----- Yellow accent line -----
    add_rect(s, Inches(0.5), Inches(1.50), Inches(12.3), Inches(0.04), YELLOW)

    # ===== LEFT: Badge box + Code rectangle =====
    badge_x = Inches(0.5); badge_y = Inches(1.72)
    badge_w = Inches(2.6); badge_h = Inches(2.4)
    add_rect(s, badge_x, badge_y, badge_w, badge_h, GRAY_BG,
             line=RGBColor(0xE5,0xE7,0xEB))
    pad = Inches(0.15)
    add_picture_fitted(s, asset_path(p.get('badge')),
                       badge_x + pad, badge_y + pad,
                       badge_w - 2*pad, badge_h - 2*pad)
    # Code rectangle (square, no outline, "Code: X")
    if p.get('code'):
        add_rect(s, badge_x, badge_y + badge_h + Inches(0.12),
                 badge_w, Inches(0.42), accent,
                 text=f'Code: {p["code"]}', size=13, bold=True, color=WHITE,
                 align='center', anchor='middle')

    # ===== RIGHT (upper): Description box =====
    rx = Inches(3.35); rw = Inches(9.45)
    # Description header bar (teal)
    add_rect(s, rx, Inches(1.72), rw, Inches(0.32), TEAL_VDARK,
             text='DESCRIPTION', size=10, bold=True, color=WHITE,
             align='left', anchor='middle', margin=0.18)
    # Description body (bigger font, in light gray box)
    desc = p.get('description') or p.get('tagline','')
    add_rect(s, rx, Inches(2.07), rw, Inches(1.15), GRAY_BG,
             text=desc, size=12, color=TEXT, anchor='top', margin=0.15)

    # ===== RIGHT (lower): What's Included =====
    # Header
    add_rect(s, rx, Inches(3.30), rw, Inches(0.32), TEAL_VDARK,
             text="WHAT'S INCLUDED", size=10, bold=True, color=WHITE,
             align='left', anchor='middle', margin=0.18)
    # Each include as its own chip rectangle
    includes = standard_includes(p)
    chip_y = Inches(3.67); chip_h = Inches(0.25); chip_gap = Inches(0.04)
    for inc in includes:
        add_rect(s, rx, chip_y, rw, chip_h, GRAY_BG,
                 text=inc, size=10.5, color=TEXT, align='left',
                 anchor='middle', margin=0.18)
        chip_y += chip_h + chip_gap

    # ===== 3 outcome cards (chip-style items) =====
    objectives, skills, roles = get_outcomes(p)
    card_y = Inches(5.10); card_h = Inches(1.35)
    margin_x = Inches(0.5); gap = Inches(0.2)
    total_w = W - 2*margin_x
    card_w = (total_w - 2*gap) // 3
    card_data = [
        ("WHAT YOU'LL LEARN", objectives, TEAL),
        ("SKILLS YOU'LL GAIN", skills,     ORANGE),
        ("POTENTIAL ROLES",    roles,      PINK),
    ]
    for i, (title, items, color) in enumerate(card_data):
        x = margin_x + i*(card_w + gap)
        # Outer card body (light bg)
        add_rect(s, x, card_y, card_w, card_h, GRAY_BG,
                 line=RGBColor(0xE5,0xE7,0xEB))
        # Header strip (colored)
        add_rect(s, x, card_y, card_w, Inches(0.30), color,
                 text=title, size=9.5, bold=True, color=WHITE,
                 align='center', anchor='middle')
        # Each item as a chip inside the card
        item_x = x + Inches(0.08); item_w = card_w - Inches(0.16)
        item_y = card_y + Inches(0.36); item_h = Inches(0.22); item_gap = Inches(0.03)
        for it in items[:4]:
            add_rect(s, item_x, item_y, item_w, item_h, WHITE,
                     text=it, size=9.5, color=TEXT, align='left',
                     anchor='middle', margin=0.08)
            item_y += item_h + item_gap

    # ===== Metadata strip (4 navy cells) =====
    cell_y = Inches(6.55); cell_h = Inches(0.5)
    cell_w = (total_w - 3*Inches(0.15)) // 4
    pairs = [
        ('FEE',      p.get('priceDisplay') or '—'),
        ('DURATION', p.get('duration') or '—'),
        ('MODALITY', (p.get('modality') or '—').split('(')[0].strip()),
        ('FOR',      p.get('target') or '—'),
    ]
    for i,(k,v) in enumerate(pairs):
        x = margin_x + i*(cell_w + Inches(0.15))
        cell = add_rect(s, x, cell_y, cell_w, cell_h, NAVY)
        tf = cell.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
        tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.LEFT
        r = p1.add_run(); r.text = k
        r.font.size = Pt(8); r.font.bold = True; r.font.name = FONT
        r.font.color.rgb = YELLOW
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run(); r2.text = v
        r2.font.size = Pt(11); r2.font.bold = True; r2.font.name = FONT
        r2.font.color.rgb = WHITE

    footer(s, page, total, section_label)

# ============ BOOKS CONSOLIDATED ============
def slide_books_grid(title, products_list, page, total):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, W, H, WHITE)
    add_rect(s, 0, 0, W, Inches(0.45), TEAL_VDARK,
             text=f'GIMI  /  Books & Guides  /  {title}',
             size=11, color=WHITE, align='left', anchor='middle', margin=0.4)
    add_rect(s, Inches(0.5), Inches(0.6), Inches(12.3), Inches(0.7), WHITE,
             text=title, size=26, bold=True, color=NAVY, anchor='middle')
    add_rect(s, Inches(0.5), Inches(1.78), Inches(12.3), Inches(0.04), YELLOW)

    n = len(products_list)
    # Choose grid columns
    if n <= 3: cols = n; rows = 1
    elif n == 4: cols = 4; rows = 1
    else: cols = (n+1)//2; rows = 2
    if cols > 5: cols = 5; rows = (n + cols - 1) // cols

    margin = Inches(0.5); gap = Inches(0.3)
    avail_w = W - 2*margin - (cols-1)*gap
    cell_w = avail_w // cols
    avail_h = Inches(4.8) - (rows-1)*gap
    cell_h = avail_h // rows
    base_y = Inches(2.0)

    for i, b in enumerate(products_list):
        row = i // cols; col = i % cols
        x = margin + col*(cell_w + gap)
        y = base_y + row*(cell_h + gap)
        # Card body
        add_rect(s, x, y, cell_w, cell_h, GRAY_BG, line=RGBColor(0xE5,0xE7,0xEB))
        # Image area (top portion)
        img_h = int(cell_h * 0.55)
        bp = asset_path(b.get('badge'))
        if bp:
            pad = Inches(0.15)
            add_picture_fitted(s, bp, x + pad, y + pad,
                               cell_w - 2*pad, img_h - pad)
        # Title strip (bigger, navy text)
        add_rect(s, x + Inches(0.1), y + img_h + Inches(0.05),
                 cell_w - Inches(0.2), Inches(0.45), GRAY_BG,
                 text=b['title'], size=13, bold=True, color=NAVY,
                 align='center', anchor='middle')
        # Tagline / short description
        if b.get('tagline'):
            add_rect(s, x + Inches(0.1), y + img_h + Inches(0.55),
                     cell_w - Inches(0.2), Inches(0.55), GRAY_BG,
                     text=b['tagline'][:90], size=10.5, color=TEXT,
                     align='center', anchor='top')
        # Price chip (navy fill, white text, prominent)
        price_h = Inches(0.45)
        add_rect(s, x + Inches(0.3), y + cell_h - price_h - Inches(0.15),
                 cell_w - Inches(0.6), price_h, NAVY,
                 text=b.get('priceDisplay','TBD'), size=16, bold=True, color=WHITE,
                 align='center', anchor='middle')
    footer(s, page, total, 'Books & Guides')

# ============ PROGRAMS / MEMBERSHIPS CONSOLIDATED ============
def slide_grid(title, products_list, page, total, section_label, accent, sub=None):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, W, H, WHITE)
    add_rect(s, 0, 0, W, Inches(0.45), TEAL_VDARK,
             text=f'GIMI  /  {section_label}  /  {title}',
             size=11, color=WHITE, align='left', anchor='middle', margin=0.4)
    add_rect(s, Inches(0.5), Inches(0.6), Inches(12.3), Inches(0.65), WHITE,
             text=title, size=24, bold=True, color=NAVY, anchor='middle')
    if sub:
        add_rect(s, Inches(0.5), Inches(1.25), Inches(12.3), Inches(0.45), WHITE,
                 text=sub, size=13, color=GRAY, italic=True, anchor='middle')
    add_rect(s, Inches(0.5), Inches(1.78), Inches(12.3), Inches(0.04), YELLOW)

    n = len(products_list)
    if n <= 4: cols = n; rows = 1
    elif n <= 8: cols = 4; rows = 2
    else: cols = 4; rows = 3
    margin = Inches(0.5); gap = Inches(0.3)
    avail_w = W - 2*margin - (cols-1)*gap
    cell_w = avail_w // cols
    base_y = Inches(2.0)
    avail_h = Inches(4.8) - (rows-1)*gap
    cell_h = avail_h // rows

    for i, item in enumerate(products_list):
        row = i // cols; col = i % cols
        x = margin + col*(cell_w + gap)
        y = base_y + row*(cell_h + gap)
        add_rect(s, x, y, cell_w, cell_h, GRAY_BG, line=RGBColor(0xE5,0xE7,0xEB))
        # Image area
        img_h = int(cell_h * 0.48)
        bp = asset_path(item.get('badge'))
        if bp:
            pad = Inches(0.15)
            add_picture_fitted(s, bp, x + pad, y + pad,
                               cell_w - 2*pad, img_h - pad)
        # Title (bigger)
        add_rect(s, x + Inches(0.1), y + img_h + Inches(0.05),
                 cell_w - Inches(0.2), Inches(0.5), GRAY_BG,
                 text=item['title'], size=12, bold=True, color=NAVY,
                 align='center', anchor='middle')
        # Tagline / description
        if item.get('tagline'):
            add_rect(s, x + Inches(0.1), y + img_h + Inches(0.57),
                     cell_w - Inches(0.2), Inches(0.6), GRAY_BG,
                     text=item['tagline'][:100], size=10, color=TEXT,
                     align='center', anchor='top')
        # Price chip (navy fill, white text)
        price_h = Inches(0.45)
        add_rect(s, x + Inches(0.3), y + cell_h - price_h - Inches(0.15),
                 cell_w - Inches(0.6), price_h, NAVY,
                 text=item.get('priceDisplay','TBD'), size=15, bold=True, color=WHITE,
                 align='center', anchor='middle')
    footer(s, page, total, section_label)

# ============ CLOSING ============
def slide_closing(page, total):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, W, H, NAVY)
    add_rect(s, 0, 0, Inches(0.4), H, TEAL)
    s.shapes.add_picture(LOGO_WHITE, W - Inches(2.6), Inches(0.5),
                         width=Inches(2.0), height=Inches(0.66))
    add_rect(s, Inches(1.0), Inches(2.2), Inches(11), Inches(1.3), NAVY,
             text='Thank you', size=72, bold=True, color=WHITE,
             align='left', anchor='middle')
    add_rect(s, Inches(1.0), Inches(3.7), Inches(11), Inches(0.6), NAVY,
             text='Explore the full catalogue at giminstitute.org/store',
             size=22, color=YELLOW, align='left')
    add_rect(s, Inches(1.0), Inches(4.7), Inches(11), Inches(0.4), NAVY,
             text='Global Innovation Management Institute',
             size=16, bold=True, color=WHITE, align='left')
    add_rect(s, Inches(1.0), Inches(5.15), Inches(11), Inches(0.4), NAVY,
             text='110 Cambridge Street, Cambridge MA 02141',
             size=12, color=RGBColor(0xCC,0xCC,0xCC), align='left')
    add_rect(s, Inches(1.0), Inches(5.55), Inches(11), Inches(0.4), NAVY,
             text='info@giminstitute.org',
             size=12, color=RGBColor(0xCC,0xCC,0xCC), align='left')

# ============ BUILD SEQUENCE ============
TOTAL = 50  # set later but final count must match

# Predefine consolidated section content
BOOK_SLIDES = [
    ('IMBOK & Future Foresight BOK',
     ['imbok-1','imbok-3','ffbok-1','ffbok-2']),
    ('Innovation Stories',
     ['game-changers','connectivate','healthovate','greenovate','squiggly']),
    ('Practitioner Tools, ISO Standards & Journal',
     ['idex','technovate','iso-book','ijis']),
]

PROGRAM_SLIDES = [
    ('Individual Assessments & Recognition',
     'For individual professionals',
     ['ipa-individual','mindset-index','awards-individual']),
    ('Organisational Assessments',
     'For organisations and business units',
     ['ipa-corporate','oia-self','full-oia']),
    ('Software & Digital Tools',
     'Platforms and software offerings',
     ['akaio','eureka','service-tool','quickgrowth']),
    ('Programs & Engagements',
     'Capability building and engagement offerings',
     ['academy','coaching','awards-org','speech']),
]

# 1 cover + 1 minto + 8 dividers + 35 cert slides + 3 book slides + 4 program slides + 1 mem slide + 1 closing
# = 1+1+8+(10+10+5+2+4)+3+4+1+1 = 50
TOTAL = 50

# Build
slide_cover()                                 # 1
slide_minto()                                 # 2
page = 3

section_order = [('icore',10), ('ielec',10), ('ff',5), ('lead',2), ('cons',4),
                 ('books',3),  ('tools',4),  ('memb',1)]

for sec_num, (sec, _) in enumerate(section_order, start=1):
    track_info = next(t for t in TRACKS if t[0]==sec)
    _, section_label, sub, accent = track_info
    sec_products = [p for p in products if p['track']==sec]

    # Section count for divider
    if sec == 'books':   sec_count = 13
    elif sec == 'memb':  sec_count = 7
    else:                sec_count = len(sec_products)

    slide_section_divider(sec_num, section_label, sub, sec_count, accent, page, TOTAL)
    page += 1

    if sec == 'books':
        for title, ids in BOOK_SLIDES:
            items = [PB[i] for i in ids if i in PB]
            slide_books_grid(title, items, page, TOTAL)
            page += 1
    elif sec == 'tools':
        for title, sub2, ids in PROGRAM_SLIDES:
            items = [PB[i] for i in ids if i in PB]
            slide_grid(title, items, page, TOTAL, 'Programs & Tools', accent, sub=sub2)
            page += 1
    elif sec == 'memb':
        slide_grid('All Memberships', sec_products, page, TOTAL,
                   'Memberships', accent,
                   sub='Annual access tiers for individuals, companies, and academia')
        page += 1
    else:
        # Cluster-order for cert sections
        if sec == 'icore':
            order = ['cip','cip-mc','ccio','ccio-mc','impact','icore-train']
        elif sec == 'ff':
            order = ['ff-series','ff-train']
        else:
            order = []
        if order:
            sorted_products = []
            for cl in order:
                sorted_products.extend([p for p in sec_products if p.get('cluster')==cl])
            sorted_products.extend([p for p in sec_products if not p.get('cluster')])
            sec_products = sorted_products
        for p in sec_products:
            slide_product(p, page, TOTAL, section_label, accent)
            page += 1

slide_closing(page, TOTAL)

# Save
OUT = r'C:\Users\aelha\Downloads\GIMI_Product_Catalog_V55.pptx'
prs.save(OUT)
print(f'Built {len(prs.slides)} slides -> {OUT}')
